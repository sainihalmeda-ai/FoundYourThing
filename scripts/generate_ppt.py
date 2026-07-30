"""Generate visual Idea Quest 2026 PPT — minimal text, diagrams, animations."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parents[1] / "docs" / "FoundYourThing_IdeaQuest2026.pptx"

# Brand palette
PRIMARY = RGBColor(0x1B, 0x4D, 0x89)
PRIMARY_LIGHT = RGBColor(0xE6, 0xF0, 0xFF)
ACCENT = RGBColor(0x0E, 0xA5, 0xA4)
ACCENT_SOFT = RGBColor(0xD6, 0xF6, 0xF5)
DARK = RGBColor(0x10, 0x2A, 0x43)
MUTED = RGBColor(0x62, 0x7D, 0x98)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFB)
ORANGE = RGBColor(0xF5, 0x9F, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

_EFFECTS = {
    "fade": ("10", "fade"),
    "appear": ("1", "appear"),
    "fly": ("2", "fly"),
}


class Build:
    """Click-to-reveal animation groups (fade on click)."""

    def __init__(self, slide):
        self.slide = slide
        self.steps: list[list[int]] = []
        self._current: list[int] = []

    def track(self, shape) -> None:
        self._current.append(shape.shape_id)

    def step(self):
        return _BuildStep(self)

    def apply(self, effect: str = "fade", duration: float = 0.6) -> None:
        if not self.steps:
            return
        preset_id, filt = _EFFECTS.get(effect, _EFFECTS["fade"])
        dur_ms = int(duration * 1000)
        parts = []
        for spids in self.steps:
            effects = []
            for i, spid in enumerate(spids):
                node = "clickEffect" if i == 0 else "withEffect"
                delay = 0 if i == 0 else 200
                effects.append(
                    f'<p:par><p:cTn id="{100 + spid}" dur="indefinite" nodeType="{node}">'
                    f'<p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>'
                    f'<p:childTnLst><p:animEffect transition="in" filter="{filt}">'
                    f'<p:cBhvr><p:cTn dur="{dur_ms}" fill="hold"/>'
                    f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>'
                    f"</p:animEffect></p:childTnLst></p:cTn></p:par>"
                )
            parts.append(
                "<p:seq concurrent=\"1\" nextAc=\"seek\">"
                "<p:cTn id=\"1\" dur=\"indefinite\" nodeType=\"mainSeq\">"
                "<p:childTnLst>" + "".join(effects) + "</p:childTnLst></p:cTn></p:seq>"
            )
        xml = (
            '<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            "<p:tnLst><p:par><p:cTn id=\"0\" dur=\"indefinite\" restart=\"never\" nodeType=\"tmRoot\">"
            "<p:childTnLst>" + "".join(parts) + "</p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"
        )
        self.slide._element.append(parse_xml(xml))


class _BuildStep:
    def __init__(self, build: Build):
        self.build = build

    def __enter__(self):
        self.build._current = []
        return self

    def __exit__(self, *_):
        if self.build._current:
            self.build.steps.append(self.build._current)
        self.build._current = []


def add_slide_transition(slide, kind: str = "fade", duration_ms: int = 800) -> None:
    xml = f"""
    <p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med">
      <p:{kind}/>
    </p:transition>
    """
    slide.element.insert(-1, parse_xml(xml))


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return slide


def bg(slide, color: RGBColor) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def accent_bar(slide, top=0, height=0.12) -> None:
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(top), SLIDE_W, Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def rounded_card(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def text_in(shape, text: str, size=16, bold=False, color=DARK, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return shape


def slide_heading(slide, title: str, subtitle: str = "") -> None:
    accent_bar(slide)
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.35), Inches(12), Inches(1.1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = MUTED


def icon_card(slide, left, top, w, h, emoji, title, caption, fill=WHITE):
    card = rounded_card(slide, left, top, w, h, fill, PRIMARY_LIGHT)
    t = slide.shapes.add_textbox(Inches(left), Inches(top + 0.15), Inches(w), Inches(0.55))
    p = t.text_frame.paragraphs[0]
    p.text = emoji
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    t2 = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.65), Inches(w - 0.2), Inches(0.55))
    p2 = t2.text_frame.paragraphs[0]
    p2.text = title
    p2.alignment = PP_ALIGN.CENTER
    p2.font.bold = True
    p2.font.size = Pt(15)
    p2.font.color.rgb = PRIMARY
    t3 = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 1.15), Inches(w - 0.2), Inches(h - 1.3))
    tf = t3.text_frame
    tf.word_wrap = True
    p3 = tf.paragraphs[0]
    p3.text = caption
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(12)
    p3.font.color.rgb = MUTED
    return card


def arrow_right(slide, left, top, width=0.55, height=0.25):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def slide_title(prs: Presentation) -> None:
    slide = blank_slide(prs)
    bg(slide, PRIMARY)
    accent_bar(slide, 6.9, 0.6)

    # Decorative circles
    for x, y, s, alpha in [(10.5, 0.8, 2.2, PRIMARY_LIGHT), (0.5, 4.5, 1.8, ACCENT_SOFT)]:
        c = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(s), Inches(s))
        c.fill.solid()
        c.fill.fore_color.rgb = alpha
        c.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = "FOUND YOUR THING"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(1.2), Inches(2.7), Inches(11), Inches(1.4))
    tf = sub.text_frame
    for i, line in enumerate(
        [
            "AI-Powered Privacy-First Campus Lost & Found",
            "(Lost on campus. Found by intelligence.)",
        ]
    ):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(20 if i == 0 else 16)
        para.font.color.rgb = WHITE if i == 0 else ACCENT_SOFT

    meta = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.3), Inches(1.8))
    tf = meta.text_frame
    for i, line in enumerate(
        [
            "Meda Sai Nihal  ·  CSE  ·  AI & Data Science",
            "SDG 11  ·  Idea Quest 2026",
        ]
    ):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(15)
        para.font.color.rgb = WHITE

    add_slide_transition(slide, "fade")


def slide_problem(prs: Presentation) -> None:
    slide = blank_slide(prs)
    bg(slide, LIGHT_BG)
    slide_heading(slide, "Problem Statement", "Campus lost & found is broken today")
    b = Build(slide)

    cards = [
        ("📱", "Lost Daily", "Phones, watches, IDs go missing"),
        ("🐌", "Too Slow", "WhatsApp & notice boards fail"),
        ("🔓", "No Privacy", "Phone numbers shared publicly"),
        ("🗑️", "Noise", "Pens & pencils clutter the system"),
    ]
    positions = [(0.7, 1.7), (3.55, 1.7), (6.4, 1.7), (9.25, 1.7)]
    with b.step():
        for (emoji, title, cap), (x, y) in zip(cards, positions):
            s = icon_card(slide, x, y, 2.75, 2.1, emoji, title, cap)
            b.track(s)

    banner = rounded_card(slide, 2.2, 4.35, 8.9, 0.95, ORANGE)
    text_in(banner, "Students need a faster, smarter & safer way to recover valuables", 18, True, WHITE)
    with b.step():
        b.track(banner)

    b.apply("fade")
    add_slide_transition(slide, "push")


def slide_solution(prs: Presentation) -> None:
    slide = blank_slide(prs)
    bg(slide, WHITE)
    slide_heading(slide, "Proposed Solution", "FoundYourThing — AI + Privacy")
    b = Build(slide)

    center = rounded_card(slide, 4.9, 2.0, 3.5, 1.5, PRIMARY)
    text_in(center, "FoundYourThing\nAI Match Engine", 18, True, WHITE)
    with b.step():
        b.track(center)

    orbit = [
        (1.0, 2.0, "📸", "Photo\nReport"),
        (10.0, 2.0, "🤖", "AI\nMatching"),
        (1.0, 4.6, "🪪", "VTU ID\nOnly"),
        (10.0, 4.6, "✅", "Mutual\nConsent"),
    ]
    with b.step():
        for x, y, em, lbl in orbit:
            c = rounded_card(slide, x, y, 2.3, 1.55, ACCENT_SOFT, ACCENT)
            text_in(c, f"{em}\n{lbl}", 16, True, DARK)
            b.track(c)

    # Connectors
    for sx, sy in [(3.3, 2.7), (8.0, 2.7), (3.3, 5.0), (8.0, 5.0)]:
        slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(sx), Inches(sy), Inches(5.0), Inches(2.75)
        ).line.color.rgb = MUTED

    tag = rounded_card(slide, 3.5, 6.05, 6.3, 0.65, PRIMARY_LIGHT)
    text_in(tag, "Valuables only  ·  No pens or pencils", 14, True, PRIMARY)
    with b.step():
        b.track(tag)

    b.apply("fade")
    add_slide_transition(slide, "fade")


def slide_architecture(prs: Presentation) -> None:
    slide = blank_slide(prs)
    bg(slide, LIGHT_BG)
    slide_heading(slide, "Architecture", "Visual system design — not just text")
    b = Build(slide)

    layers = [
        (0.8, 2.0, 3.5, "📲 Mobile App", "Expo · Android\niOS · Web", PRIMARY, WHITE),
        (4.9, 2.0, 3.5, "⚡ API Layer", "FastAPI · Auth\nItems · Claims", ACCENT, WHITE),
        (9.0, 2.0, 3.5, "🧠 AI + Data", "Embeddings\nSQLite / DB", DARK, WHITE),
    ]
    with b.step():
        for x, y, w, title, sub, fill, txt in layers:
            box = rounded_card(slide, x, y, w, 2.2, fill)
            text_in(box, f"{title}\n\n{sub}", 15, True, txt)
            b.track(box)

    with b.step():
        arrow_right(slide, 4.35, 2.95)
        arrow_right(slide, 8.45, 2.95)

    # Bottom services row
    services = ["Image AI", "Text AI", "Privacy", "Uploads"]
    with b.step():
        for i, name in enumerate(services):
            s = rounded_card(slide, 1.0 + i * 2.9, 4.8, 2.5, 0.85, WHITE, ACCENT)
            text_in(s, name, 14, True, PRIMARY)
            b.track(s)

    # User flow strip
    flow_labels = ["Report", "Match", "Claim", "Connect"]
    with b.step():
        for i, lbl in enumerate(flow_labels):
            x = 1.3 + i * 2.85
            c = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(6.0), Inches(0.75), Inches(0.75)
            )
            c.fill.solid()
            c.fill.fore_color.rgb = ACCENT if i < 3 else PRIMARY
            c.line.fill.background()
            text_in(c, str(i + 1), 16, True, WHITE)
            t = slide.shapes.add_textbox(Inches(x - 0.35), Inches(6.85), Inches(1.45), Inches(0.35))
            p = t.text_frame.paragraphs[0]
            p.text = lbl
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
            p.font.color.rgb = DARK
            b.track(c)
            b.track(t)
            if i < 3:
                arrow_right(slide, x + 0.85, 6.25, 0.45, 0.2)

    b.apply("fade")
    add_slide_transition(slide, "wipe")


def slide_impact(prs: Presentation) -> None:
    slide = blank_slide(prs)
    bg(slide, WHITE)
    slide_heading(slide, "Expected Impact", "Why this matters on campus")
    b = Build(slide)

    stats = [
        ("⚡", "Faster Recovery", "AI match\nin seconds"),
        ("🔐", "Safe Contact", "Phone after\nconsent"),
        ("🏫", "SDG 11", "Better campus\ncommunity"),
        ("♻️", "Less Waste", "Fewer replaced\nvaluables"),
    ]
    with b.step():
        for i, (em, title, val) in enumerate(stats):
            x = 0.85 + i * 3.05
            card = rounded_card(slide, x, 2.0, 2.75, 3.2, PRIMARY_LIGHT if i % 2 == 0 else ACCENT_SOFT)
            t = slide.shapes.add_textbox(Inches(x), Inches(2.15), Inches(2.75), Inches(0.6))
            p = t.text_frame.paragraphs[0]
            p.text = em
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(34)
            text_in(card, f"{title}\n\n{val}", 16, True, PRIMARY if i % 2 == 0 else DARK)
            b.track(card)

    b.apply("fly")
    add_slide_transition(slide, "fade")


def slide_feasibility(prs: Presentation) -> None:
    slide = blank_slide(prs)
    bg(slide, LIGHT_BG)
    slide_heading(slide, "Feasibility & Plan", "Phased rollout")
    b = Build(slide)

    phases = [
        ("Phase 1", "Core App", "Auth · Reports\nAI MVP", "✅"),
        ("Phase 2", "Better AI", "CLIP · Accuracy", "⏳"),
        ("Phase 3", "Campus", "Alerts · VTU API", "⏳"),
        ("Phase 4", "Scale", "Admin · Cloud", "⏳"),
    ]
    with b.step():
        for i, (phase, title, detail, status) in enumerate(phases):
            x = 0.75 + i * 3.05
            card = rounded_card(slide, x, 2.0, 2.75, 3.5, WHITE, PRIMARY_LIGHT)
            badge = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.95), Inches(2.15), Inches(0.85), Inches(0.85)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = ACCENT if status == "✅" else MUTED
            badge.line.fill.background()
            text_in(badge, status, 18, True, WHITE)
            t = slide.shapes.add_textbox(Inches(x), Inches(3.1), Inches(2.75), Inches(2.0))
            tf = t.text_frame
            for j, line in enumerate([phase, title, detail]):
                para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                para.text = line
                para.alignment = PP_ALIGN.CENTER
                para.font.size = Pt(13 if j == 2 else 15)
                para.font.bold = j < 2
                para.font.color.rgb = PRIMARY if j == 1 else MUTED if j == 0 else DARK
            b.track(card)
            if i < 3:
                arrow_right(slide, x + 2.85, 3.55, 0.35, 0.22)

    b.apply("fade")
    add_slide_transition(slide, "push")


def slide_future(prs: Presentation) -> None:
    slide = blank_slide(prs)
    bg(slide, WHITE)
    slide_heading(slide, "Future Scope", "Where FoundYourThing can grow")
    b = Build(slide)

    items = [
        "🔔 Push alerts (valuables only)",
        "🏛️ College SSO integration",
        "📊 Admin analytics dashboard",
        "🌐 Multi-campus support",
        "🏷️ QR tags on items",
        "🛡️ Auto-blur ID card photos",
    ]
    with b.step():
        for i, item in enumerate(items):
            col = i % 2
            row = i // 2
            x = 1.0 + col * 6.0
            y = 1.75 + row * 1.45
            pill = rounded_card(slide, x, y, 5.5, 0.95, ACCENT_SOFT if col else PRIMARY_LIGHT, ACCENT)
            text_in(pill, item, 16, True, DARK, PP_ALIGN.LEFT)
            b.track(pill)

    b.apply("fade")
    add_slide_transition(slide, "fade")


def slide_demo(prs: Presentation) -> None:
    slide = blank_slide(prs)
    bg(slide, PRIMARY)
    slide_heading(slide, "Demo / Prototype", "Working application today")
    slide.shapes[-1].text_frame.paragraphs[0].font.color.rgb = WHITE
    slide.shapes[-1].text_frame.paragraphs[1].font.color.rgb = ACCENT_SOFT

    features = ["Login", "Report", "Feed", "AI Match", "Claim"]
    b = Build(slide)
    with b.step():
        for i, f in enumerate(features):
            x = 0.9 + i * 2.45
            c = rounded_card(slide, x, 2.3, 2.15, 1.2, WHITE)
            text_in(c, f, 15, True, PRIMARY)
            b.track(c)
            if i < 4:
                arrow_right(slide, x + 2.2, 2.75, 0.35, 0.22)

    stack = rounded_card(slide, 3.8, 4.0, 5.7, 1.5, ACCENT)
    text_in(stack, "Expo Mobile  +  Web  +  Python FastAPI", 18, True, WHITE)
    with b.step():
        b.track(stack)

    note = slide.shapes.add_textbox(Inches(2.5), Inches(5.85), Inches(8.3), Inches(0.6))
    p = note.text_frame.paragraphs[0]
    p.text = "Next slides → Add your prototype screenshots"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE

    b.apply("fade")
    add_slide_transition(slide, "fade")


def slide_photo_placeholder(prs: Presentation, num: int) -> None:
    slide = blank_slide(prs)
    bg(slide, LIGHT_BG)
    slide_heading(slide, f"Prototype Screenshot {num - 8}", "Insert your app image here")

    frame = rounded_card(slide, 1.5, 1.65, 10.3, 4.85, WHITE, ACCENT)
    inner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.85), Inches(2.0), Inches(9.6), Inches(4.15)
    )
    inner.fill.solid()
    inner.fill.fore_color.rgb = PRIMARY_LIGHT
    inner.line.color.rgb = ACCENT
    inner.line.width = Pt(2)
    inner.line.dash_style = 2  # dashed

    label = slide.shapes.add_textbox(Inches(2.5), Inches(3.6), Inches(8.3), Inches(1.2))
    p = label.text_frame.paragraphs[0]
    p.text = "[ ADD PHOTO HERE ]"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = MUTED

    hints = {
        9: "Home screen · Report Lost · Login",
        10: "AI Match · Claim flow · Contact share",
    }
    hint = slide.shapes.add_textbox(Inches(2.5), Inches(6.55), Inches(8.3), Inches(0.4))
    p2 = hint.text_frame.paragraphs[0]
    p2.text = hints.get(num, "")
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(13)
    p2.font.color.rgb = MUTED

    add_slide_transition(slide, "fade")


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_impact(prs)
    slide_feasibility(prs)
    slide_future(prs)
    slide_demo(prs)
    slide_photo_placeholder(prs, 9)
    slide_photo_placeholder(prs, 10)

    return prs


def main() -> None:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    build().save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
